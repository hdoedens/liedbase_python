from __future__ import print_function
from pptx import Presentation
import argparse
import common
import re
import helpers

bijbelboeken = ['genesis']
liedboeken = ['gezang','levenslied','lied','opwekking','psalm']

LITURGIE_FILENAME = "liturgie.txt"
TEMPLATE_FILENAME = "template.pptx"
TEMPLATE_TEXT_PLACEHOLDER = "Text Placeholder 1"
OUTPUT_FILENAME = 'presentatie.pptx'
MAX_NUM_LINES_PER_SHEET = 4

def parse_args():
    """ Setup the input and output arguments for the script
    Return the parsed input and output files
    """
    parser = argparse.ArgumentParser(description='Analyze powerpoint file structure and create chuch presentation with it')
    parser.add_argument('analyze',
                        type=str,
                        nargs='?',
                        help='"analyze" voor analyze presentatie template')
    return parser.parse_args()

def readInputFile(inputFile):
    f = open(inputFile, 'r')
    # loop over lines to filter empty lines
    for line in f.readlines():
        # remove unwanted spaces and newlines
        line = line.strip()
        # treat lines starting with hash-sign as comment
        if not line or line[0:1] == '#':
            continue
        
        #validate syntax
        if(helpers.validate_line(line.lower())):
            getLiturgyLineContent(line)

def getLiturgyLineContent(line):
    print("Slides maken voor: {}".format(line))
    source = line[0:line.find(' ')]
    
    if helpers.get_bijbelboek(line) != None:
        bijbeltekst = get_bijbeltekst_from_source(line)
        maak_slides(bijbeltekst)
        
    elif source in liedboeken:
        source_path = "bronnen/liederen/{}.txt".format(source)
        print("Als bron wordt gebruikt: {}".format(source_path))
        # haal liedtekst op
        f = open(source_path, 'r')
        
        # alle liederen met verzen 
        if(source != "opwekking"):
            liederen = get_lines_from_source(line, source, f)
        else:
            # opwekking, zonder verzen
            liederen = get_opwekking_from_source(line, source, f)
        
        maak_slides(liederen)
        
    else:
        print("Dit type slide wordt (nog) niet ondersteund")

def get_lines_from_source(line, source, f):
    # uitgangspunt: verzen worden altijd meegegeven als komma-gescheiden lijst.
        # voorbeeld:
        #   psalm 4: 1
        #   gezang 14: 3, 6
    songNumber = line[line.find(' '):line.find(':')].strip()
    verses = line[line.find(':')+1:].split(',')
    zoekstring = "^{} {}:".format(source, songNumber)
    print("De volgende zoekstring wordt gebruikt: {}".format(zoekstring))
    liedGevonden = False
    versGevonden = False
    currentVerse = -1
    
    curr_num_lines = 0
        
        # gevonden regels opslaan in 2d array
    liederen = {}
        # for verse in verses:
        #     liederen[verse] = []
            
    for regel in f.readlines():
        regel = regel.strip()
            # zoek het juiste lied op
        if(re.search(zoekstring, regel)):
            liedGevonden = True
            continue
            # bij het volgende lied stoppen
        if(re.search("^{} ".format(source), regel)):
            if liedGevonden:
                break
        if liedGevonden:
            if(regel in verses):
                currentVerse = regel
                versGevonden = True
                liederen[currentVerse] = []
                curr_num_lines = 0
                continue
            elif("" == regel or (common.is_integer(regel) and regel not in verses)):
                versGevonden = False
            if(versGevonden):
                if curr_num_lines == MAX_NUM_LINES_PER_SHEET:
                    currentVerse = currentVerse + 'a'
                    curr_num_lines = 0
                        # add the new subverse to the dictionary
                    liederen[currentVerse] = []
                liederen[currentVerse].append(regel)
                curr_num_lines += 1
    return liederen

def get_opwekking_from_source(line, source, f):
    # uitgangspunt: opwekkingsliederen kennen geen verzen
    # dus het hele lied wordt altijd gekopieerd.
    songNumber = line[line.find(' '):].strip()
    zoekstring = "^{} {}".format(source, songNumber)
    print("De volgende zoekstring wordt gebruikt: {}".format(zoekstring))
    liedGevonden = False
        
    max_num_lines = 4
    curr_num_lines = 0
        
    # gevonden regels opslaan in 2d array met virtueel en fixed vers nummer '1'
    liederen = {}
    currentVerse = 'a'
    liederen[currentVerse] = [] #initialiseren
        
    for regel in f.readlines():
        regel = regel.strip()
        # zoek het juiste lied op
        if(re.search(zoekstring, regel)):
            liedGevonden = True
            continue
        # bij het volgende lied stoppen
        if(re.search("^{} ".format(source), regel)):
            if liedGevonden:
                break
        if liedGevonden:
            if curr_num_lines == max_num_lines:
                currentVerse = currentVerse + 'a'
                curr_num_lines = 0
                # add the new subverse to the dictionary
                liederen[currentVerse] = []
            liederen[currentVerse].append(regel)
            curr_num_lines += 1
    return liederen
        
def get_bijbeltekst_from_source(line):
    source = helpers.get_source(line)
    sourcepath = "./bronnen/bijbels/BGT/{}.txt".format(source)
    chapter = helpers.get_chapter(line)
    van = helpers.get_van_vers(line)
    tot = helpers.get_tot_vers(line)
    print("Bijbeltekst slide maken met als bron: {}".format(sourcepath))
    print("Boek {} hoofdstuk {} verzen {} tot en met {}".format(source, chapter, van, tot))
    lines = helpers.get_text_from_bible("BGT", source, chapter, van, tot)
    if not lines:
        raise Exception("Geen bijbeltekst gevonden voor: {}".format(line))
    inhoud = {}
    inhoud['a'] = lines
    return inhoud
    
def maak_slides(inhoud):
    for liednummer, regels in inhoud.items():
        textslide = prs.slides.add_slide(prs.slide_layouts[0])
        for shape in textslide.placeholders:
            if shape.is_placeholder:
                phf = shape.placeholder_format
                if shape.name == TEMPLATE_TEXT_PLACEHOLDER:
                    try:
                        shape.text = "\n".join(regels)
                    except AttributeError:
                        print("{} has no text attribute".format(phf.type))

def initialize_pptx(input):
    """ Take the input file and analyze the structure.
    The output file contains marked up information to make it easier
    for generating future powerpoint templates.
    """
    global prs
    prs = Presentation(input)

def analyze_presentation(output):
    """ Take the input file and analyze the structure.
    The output file contains marked up information to make it easier
    for generating future powerpoint templates.
    """
    # Each powerpoint file has multiple layouts
    # Loop through them all and  see where the various elements are
    for index, _ in enumerate(prs.slide_layouts):
        slide = prs.slides.add_slide(prs.slide_layouts[index])
        # Not every slide has to have a title
        try:
            title = slide.shapes.title
            title.text = 'Title for Layout {}'.format(index)
        except AttributeError:
            print("No Title for Layout {}".format(index))
        # Go through all the placeholders and identify them by index and type
        for shape in slide.placeholders:
            if shape.is_placeholder:
                phf = shape.placeholder_format
                # Do not overwrite the title which is just a special placeholder
                try:
                    if 'Title' not in shape.text:
                        shape.text = '{}'.format(shape.name)
                except AttributeError:
                    print("{} has no text attribute".format(phf.type))
                print('{} {}'.format(phf.idx, shape.name))
    prs.save(output)

def save_pptx(filename):
    prs.save(filename)

if __name__ == "__main__":
    args = parse_args()
    initialize_pptx(TEMPLATE_FILENAME)
    if args.analyze == 'analyze':
        print("analyseer het template bestand")
        analyze_presentation("analyze_{}".format(OUTPUT_FILENAME))
    else:
        readInputFile(LITURGIE_FILENAME)
        save_pptx(OUTPUT_FILENAME)